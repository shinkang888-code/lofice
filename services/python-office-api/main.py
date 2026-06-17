#!/usr/bin/env python3
"""python-office 배치 API — Phase 2 (Vercel serverless 호환)"""

from __future__ import annotations

import base64
import io
import os
import tempfile
import zipfile
from pathlib import Path
from typing import Any

import uvicorn
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware

app = FastAPI(title="lofice python-office API", version="1.1.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def ppt_to_pdf_fallback(data: bytes, suffix: str) -> bytes:
    """python-pptx + reportlab — LibreOffice/COM 없이 Vercel Linux에서 동작"""
    from pptx import Presentation
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.pdfgen import canvas

    prs = Presentation(io.BytesIO(data))
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=landscape(letter))
    width, height = landscape(letter)
    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        c.setFont("Helvetica", 14)
        c.drawString(40, height - 40, f"Slide {i + 1}")
        y = height - 80
        for t in texts[:12]:
            c.drawString(40, y, t[:120])
            y -= 22
        c.showPage()
    c.save()
    return buf.getvalue()


@app.get("/health")
def health() -> dict[str, Any]:
    backends: list[str] = []
    try:
        import office  # noqa: F401
        backends.append("python-office")
    except ImportError:
        pass
    try:
        import pptx  # noqa: F401
        import reportlab  # noqa: F401
        backends.append("pptx-reportlab-fallback")
    except ImportError:
        pass
    return {
        "status": "ok",
        "backends": backends,
        "features": ["ppt-to-pdf", "excel-split", "analyze-excel"],
    }


@app.post("/convert/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)) -> dict[str, str]:
    suffix = Path(file.filename or "doc.pptx").suffix.lower()
    if suffix not in {".ppt", ".pptx", ".pptm"}:
        raise HTTPException(status_code=400, detail="only ppt/pptx/pptm")
    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="empty file")

    pdf_bytes: bytes | None = None
    try:
        from office.api import ppt

        with tempfile.TemporaryDirectory() as tmp:
            src = Path(tmp) / f"input{suffix}"
            dst = Path(tmp) / "output.pdf"
            src.write_bytes(data)
            ppt.ppt2pdf(str(src), str(dst))
            if dst.is_file():
                pdf_bytes = dst.read_bytes()
    except Exception:
        pdf_bytes = None

    if pdf_bytes is None:
        try:
            pdf_bytes = ppt_to_pdf_fallback(data, suffix)
        except Exception as e:
            raise HTTPException(status_code=500, detail=str(e)[:500]) from e

    return {
        "file_name": Path(file.filename or "slides").stem + ".pdf",
        "data_base64": base64.b64encode(pdf_bytes).decode(),
        "format": "pdf",
    }


@app.post("/convert/excel-split")
async def excel_split(file: UploadFile = File(...)) -> dict[str, str]:
    try:
        import openpyxl
    except ImportError:
        raise HTTPException(status_code=503, detail="openpyxl not installed")

    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="empty file")

    wb = openpyxl.load_workbook(io.BytesIO(data))
    out_buf = io.BytesIO()
    with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for name in wb.sheetnames:
            single = openpyxl.Workbook()
            single.remove(single.active)
            single.create_sheet(name)
            ws_src = wb[name]
            ws_dst = single[name]
            for row in ws_src.iter_rows():
                for cell in row:
                    ws_dst[cell.coordinate].value = cell.value
            sheet_buf = io.BytesIO()
            single.save(sheet_buf)
            zf.writestr(f"{name}.xlsx", sheet_buf.getvalue())
    stem = Path(file.filename or "workbook").stem
    return {
        "file_name": f"{stem}_sheets.zip",
        "data_base64": base64.b64encode(out_buf.getvalue()).decode(),
        "format": "zip",
    }


@app.post("/analyze/excel")
async def analyze_excel(file: UploadFile = File(...)) -> dict[str, Any]:
    try:
        import openpyxl
    except ImportError:
        raise HTTPException(status_code=503, detail="openpyxl not installed")

    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="empty file")

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets = []
    for name in wb.sheetnames:
        ws = wb[name]
        sheets.append({"name": name, "max_row": ws.max_row, "max_column": ws.max_column})
    wb.close()
    return {"sheet_count": len(sheets), "sheets": sheets}


if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=int(os.environ.get("PORT", "8300")))
