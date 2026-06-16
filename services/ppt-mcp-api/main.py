#!/usr/bin/env python3
"""Office-PowerPoint-MCP-Server FastAPI — python-pptx REST wrapper for lofice."""

from __future__ import annotations

import io
import json
import os
from pathlib import Path
from typing import Any, Optional

import uvicorn
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "").strip()
OPENAI_BASE = os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1").rstrip("/")
OPENAI_MODEL = os.environ.get("OPENAI_MODEL", "gpt-4o-mini")

PPT_MCP_DIR = Path(os.environ.get("PPT_MCP_DIR", Path(__file__).resolve().parent))
TEMPLATES_JSON = PPT_MCP_DIR / "slide_layout_templates.json"

app = FastAPI(title="lofice PPT MCP API", version="1.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])


def require_pptx() -> None:
    if not HAS_PPTX:
        raise HTTPException(status_code=503, detail="python-pptx not installed. pip install python-pptx")


class TemplateSequenceItem(BaseModel):
    template_id: str
    content: dict[str, str] = Field(default_factory=dict)


class CreateFromTemplatesRequest(BaseModel):
    sequence: list[TemplateSequenceItem]
    color_scheme: str = "modern_blue"
    file_name: str = "presentation.pptx"


class AutoGenerateRequest(BaseModel):
    topic: str
    slide_count: int = 6
    color_scheme: str = "modern_blue"
    language: str = "ko"


class AiDeckRequest(BaseModel):
    topic: str
    slide_count: int = 6
    language: str = "ko"
    author: str = "lofice"
    color_scheme: str = "modern_blue"


class GeneratorRequest(BaseModel):
    user_text: str
    number_of_slide: int = 6
    template_choice: str = "classic"
    presentation_title: str = ""
    presenter_name: str = "lofice"
    insert_image: bool = False
    language: str = "ko"


async def _llm_json_prompt(prompt: str) -> list[Any]:
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OPENAI_API_KEY required")
    import httpx

    async with httpx.AsyncClient(timeout=90) as client:
        r = await client.post(
            f"{OPENAI_BASE}/chat/completions",
            headers={"Authorization": f"Bearer {OPENAI_API_KEY}"},
            json={"model": OPENAI_MODEL, "messages": [{"role": "user", "content": prompt}], "temperature": 0.7},
        )
    if r.status_code != 200:
        raise HTTPException(status_code=502, detail=r.text[:500])
    content = r.json()["choices"][0]["message"]["content"]
    try:
        return json.loads(content.strip().removeprefix("```json").removesuffix("```").strip())
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=502, detail=f"LLM JSON parse error: {e}") from e


def _outline_to_sequence(outline: list[Any]) -> list[TemplateSequenceItem]:
    seq: list[TemplateSequenceItem] = []
    for item in outline:
        if not isinstance(item, dict):
            continue
        t = str(item.get("type") or item.get("slide_type") or "textual")
        title = str(item.get("title") or "슬라이드")
        content: dict[str, str] = {"title": title}
        if t == "intro":
            content["subtitle"] = str(item.get("subtitle") or "")
            seq.append(TemplateSequenceItem(template_id="title_slide", content=content))
        else:
            bullets = item.get("content") or []
            if isinstance(bullets, list):
                content["content"] = "\n".join(str(x) for x in bullets)
            else:
                content["content"] = str(bullets)
            seq.append(TemplateSequenceItem(template_id="text_with_image", content=content))
    return seq


def _extract_text_from_prs(prs: "Presentation") -> dict[str, Any]:
    slides_text = []
    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        title = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if not t:
                continue
            if not title:
                title = t
            else:
                texts.append(t)
        combined = "\n".join([title, *texts]) if title else "\n".join(texts)
        slides_text.append(
            {
                "slide_index": i,
                "slide_title": title or f"Slide {i + 1}",
                "placeholders": [{"text": title}] if title else [],
                "text_shapes": [{"text": x} for x in texts],
                "all_text_combined": combined,
                "has_notes": False,
            }
        )
    combined_all = "\n\n".join(f"=== SLIDE {s['slide_index'] + 1} ===\n{s['all_text_combined']}" for s in slides_text)
    return {
        "total_slides": len(slides_text),
        "slides_with_text": sum(1 for s in slides_text if s["all_text_combined"]),
        "slides_text": slides_text,
        "all_presentation_text_combined": combined_all,
    }


def _add_slide_with_content(prs: "Presentation", title: str, body_lines: list[str]) -> None:
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body = ph
            break
    if body is None and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
    if body is not None:
        tf = body.text_frame
        tf.clear()
        for j, line in enumerate(body_lines or [" "]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0


def _build_from_sequence(sequence: list[TemplateSequenceItem]) -> bytes:
    require_pptx()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for item in sequence:
        c = item.content
        title = c.get("title") or c.get("quote") or item.template_id.replace("_", " ").title()
        lines: list[str] = []
        for key in ("subtitle", "content", "agenda_items", "steps", "content_left", "content_right", "metric_1_value", "metric_2_value", "metric_3_value", "author", "contact"):
            if c.get(key):
                lines.append(str(c[key]))
        _add_slide_with_content(prs, title, lines)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@app.get("/health")
def health():
    tpl_count = 0
    if TEMPLATES_JSON.is_file():
        try:
            tpl_count = len(json.loads(TEMPLATES_JSON.read_text(encoding="utf-8")).get("templates", {}))
        except Exception:
            tpl_count = 0
    return {"ok": True, "python_pptx": HAS_PPTX, "templates": tpl_count, "ppt_generator": True}


@app.get("/templates")
def list_templates():
    if not TEMPLATES_JSON.is_file():
        return {"templates": [], "source": "builtin"}
    data = json.loads(TEMPLATES_JSON.read_text(encoding="utf-8"))
    ids = list(data.get("templates", {}).keys())
    return {"templates": ids, "color_schemes": list(data.get("color_schemes", {}).keys())}


@app.post("/extract-text")
async def extract_text(file: UploadFile = File(...)):
    require_pptx()
    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="empty file")
    try:
        prs = Presentation(io.BytesIO(raw))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"invalid pptx: {e}") from e
    return _extract_text_from_prs(prs)


@app.post("/create-from-templates")
def create_from_templates(req: CreateFromTemplatesRequest):
    if not req.sequence:
        raise HTTPException(status_code=400, detail="sequence required")
    from fastapi.responses import Response

    data = _build_from_sequence(req.sequence)
    return Response(content=data, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


@app.post("/auto-generate")
async def auto_generate(req: AutoGenerateRequest):
    prompt = f"""Create a presentation outline about: {req.topic}
Language: {req.language}
Return JSON array of {req.slide_count} slides: [{{"template_id":"title_slide|text_with_image|key_metrics_dashboard|thank_you_slide","content":{{"title":"...","content":"..."}}}}]
Only valid JSON, no markdown."""
    sequence_raw = await _llm_json_prompt(prompt)
    seq = [TemplateSequenceItem(**x) for x in sequence_raw]
    from fastapi.responses import Response

    return Response(content=_build_from_sequence(seq), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


@app.post("/ai-deck")
async def ai_deck(req: AiDeckRequest):
    """powerpoint Ruby gem 형식 AI outline → JSON (lofice 클라이언트가 PPTX 생성)"""
    if OPENAI_API_KEY:
        prompt = f"""Create a PowerPoint presentation outline about: {req.topic}
Language: {req.language}
Author for intro subtitle: {req.author}
Return JSON array of exactly {req.slide_count} slides using powerpoint gem types:
[{{"type":"intro","title":"Main Title","subtitle":"Subtitle"}},
 {{"type":"textual","title":"Section","content":["bullet 1","bullet 2"]}}]
Types: intro, textual (pictorial optional with image_url). JSON only, no markdown."""
        outline = await _llm_json_prompt(prompt)
    else:
        outline = [
            {"type": "intro", "title": req.topic, "subtitle": f"{req.author} · lofice"},
            {"type": "textual", "title": "개요", "content": [f"{req.topic} 소개", "핵심 논점"]},
            {"type": "textual", "title": "결론", "content": ["요약", "감사합니다"]},
        ]
    return {"outline": outline, "topic": req.topic, "source": "openai" if OPENAI_API_KEY else "builtin"}


@app.post("/ai-deck/pptx")
async def ai_deck_pptx(req: AiDeckRequest):
    """AI outline → PPTX bytes (서버 python-pptx)"""
    data = await ai_deck(req)
    seq = _outline_to_sequence(data["outline"])
    from fastapi.responses import Response

    return Response(content=_build_from_sequence(seq), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


@app.post("/generator/generate")
async def generator_generate(req: GeneratorRequest):
    """PowerPoint-Generator-Python-Project — GPT + Pexels + 테마 PPTX"""
    require_pptx()
    from urllib.parse import quote

    from ppt_generator import generate_presentation

    pptx, slides, source = await generate_presentation(
        user_text=req.user_text,
        number_of_slide=req.number_of_slide,
        template_choice=req.template_choice,
        presentation_title=req.presentation_title or req.user_text[:80],
        presenter_name=req.presenter_name,
        insert_image=req.insert_image,
        language=req.language,
    )
    from fastapi.responses import Response

    return Response(
        content=pptx,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "X-Source": source,
            "X-Slides-Json": quote(json.dumps(slides, ensure_ascii=False)),
        },
    )


@app.post("/generator/outline")
async def generator_outline(req: GeneratorRequest):
    """GPT 슬라이드 outline JSON만 반환"""
    from ppt_generator import (
        build_user_message,
        builtin_slides,
        chat_development,
        parse_response,
    )

    if OPENAI_API_KEY:
        raw = await chat_development(build_user_message(req.number_of_slide, req.user_text, req.language))
        slides = parse_response(raw)
        source = "openai"
    else:
        slides = builtin_slides(req.user_text, req.number_of_slide)
        source = "builtin"
    return {
        "slides": slides,
        "presentation_title": req.presentation_title or req.user_text[:80],
        "presenter_name": req.presenter_name,
        "template_choice": req.template_choice,
        "source": source,
    }


if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=int(os.environ.get("PORT", "8300")))
