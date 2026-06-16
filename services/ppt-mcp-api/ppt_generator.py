"""PowerPoint-Generator-Python-Project 이식 — GPT + Pexels + python-pptx."""

from __future__ import annotations

import io
import json
import os
from typing import Any
from urllib.parse import quote_plus

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "").strip()
OPENAI_BASE = os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1").rstrip("/")
OPENAI_MODEL = os.environ.get("OPENAI_MODEL", "gpt-4o-mini")
PEXELS_API_KEY = os.environ.get("PEXELS_API_KEY", "").strip()

THEME_FONTS: dict[str, dict[str, Any]] = {
    "dark_modern": {"title": "Times New Roman", "title_rgb": RGBColor(255, 165, 0), "body_rgb": RGBColor(255, 255, 255)},
    "bright_modern": {"title": "Arial", "title_rgb": RGBColor(255, 20, 147), "body_rgb": RGBColor(0, 0, 0)},
    "classic": {"title": "Arial", "title_rgb": RGBColor(0, 51, 102), "body_rgb": RGBColor(0, 0, 0)},
}

SYSTEM_PROMPT = (
    "You are an assistant that gives the idea for PowerPoint presentations. When answering, give the user "
    "the summarized content for each slide based on the number of slide. "
    "And the format of the answer must be Slide X(the number of the slide): {title of the content} /n "
    "Content: /n content with some bullet points. "
    "Keyword: /n Give the most important keyword(within two words) that represents the slide for each one"
)


def build_user_message(number_of_slide: int, user_text: str, language: str = "ko") -> str:
    return (
        f"I want you to come up with the idea for the PowerPoint. The number of slides is {number_of_slide}. "
        f"The content is: {user_text}. The title of content for each slide must be unique, "
        f"and extract the most important keyword within two words for each slide. Summarize the content for each slide. "
        f"Write in {language}."
    )


async def chat_development(user_message: str) -> str:
    if not OPENAI_API_KEY:
        raise ValueError("OPENAI_API_KEY required")
    async with httpx.AsyncClient(timeout=90) as client:
        r = await client.post(
            f"{OPENAI_BASE}/chat/completions",
            headers={"Authorization": f"Bearer {OPENAI_API_KEY}"},
            json={
                "model": OPENAI_MODEL,
                "messages": [
                    {"role": "system", "content": SYSTEM_PROMPT},
                    {"role": "user", "content": user_message},
                ],
                "temperature": 0.7,
            },
        )
    if r.status_code != 200:
        raise RuntimeError(r.text[:500])
    return r.json()["choices"][0]["message"]["content"]


def parse_response(response: str) -> list[dict[str, str]]:
    """GPT 응답 → [{title, content, keyword}, ...]"""
    slides = [s.strip() for s in response.split("\n\n") if s.strip()]
    slides_content: list[dict[str, str]] = []
    for slide in slides:
        lines = [ln for ln in slide.split("\n") if ln.strip()]
        if not lines:
            continue
        title_line = lines[0]
        if ": " in title_line:
            title = title_line.split(": ", 1)[1]
        else:
            title = title_line
        content_lines = [ln for ln in lines[1:] if ln.strip() not in ("Content:", "Contents:")]
        keyword = ""
        body_lines: list[str] = []
        for ln in content_lines:
            low = ln.lower()
            if low.startswith("keyword:") or low.startswith("keywords:"):
                keyword = ln.split(":", 1)[1].strip()
            else:
                body_lines.append(ln)
        content = "\n".join(body_lines)
        slides_content.append({"title": title, "content": content, "keyword": keyword or title.split()[0][:20]})
    return slides_content


def builtin_slides(user_text: str, number_of_slide: int) -> list[dict[str, str]]:
    """OPENAI 없을 때 fallback"""
    base = user_text.strip() or "Presentation"
    out: list[dict[str, str]] = []
    for i in range(max(1, number_of_slide)):
        out.append(
            {
                "title": f"{base} — 슬라이드 {i + 1}",
                "content": f"• {base} 관련 핵심 내용\n• lofice AI 생성 (fallback)",
                "keyword": base.split()[0][:20] if base.split() else "topic",
            }
        )
    return out


async def search_pexels_images(keyword: str) -> str | None:
    if not PEXELS_API_KEY or not keyword.strip():
        return None
    query = quote_plus(keyword.lower())
    url = f"https://api.pexels.com/v1/search?query={query}&per_page=1"
    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.get(url, headers={"Authorization": PEXELS_API_KEY})
    if r.status_code != 200:
        return None
    data = r.json()
    photos = data.get("photos") or []
    if not photos:
        return None
    return photos[0].get("src", {}).get("medium")


def _style_run(run, font_name: str, rgb: RGBColor) -> None:
    run.font.name = font_name
    run.font.color.rgb = rgb


def _style_placeholder(placeholder, theme: dict[str, Any], is_title: bool) -> None:
    rgb = theme["title_rgb"] if is_title else theme["body_rgb"]
    for paragraph in placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            _style_run(run, theme["title"], rgb)


def _add_title_slide(prs: Presentation, title: str, subtitle: str, theme_key: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    theme = THEME_FONTS.get(theme_key, THEME_FONTS["classic"])
    slide.shapes.title.text = title
    _style_placeholder(slide.shapes.title, theme, True)
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
        _style_placeholder(slide.placeholders[1], theme, False)


def _add_content_slide(
    prs: Presentation,
    slide_content: dict[str, str],
    theme_key: str,
    insert_image: bool,
    image_bytes: bytes | None,
) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    theme = THEME_FONTS.get(theme_key, THEME_FONTS["classic"])
    for ph in slide.placeholders:
        pf = ph.placeholder_format
        if pf.type == 1:
            ph.text = slide_content["title"]
            _style_placeholder(ph, theme, True)
        elif pf.type == 7:
            ph.text = slide_content["content"]
            _style_placeholder(ph, theme, False)
    if insert_image and image_bytes:
        stream = io.BytesIO(image_bytes)
        slide.shapes.add_picture(stream, Inches(8.5), Inches(3.5), width=Inches(4), height=Inches(2.5))


def create_ppt_bytes(
    slides_content: list[dict[str, str]],
    template_choice: str,
    presentation_title: str,
    presenter_name: str,
    insert_image: bool,
    image_map: dict[int, bytes | None] | None = None,
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs, presentation_title, f"Presented by {presenter_name}", template_choice)
    for i, sc in enumerate(slides_content):
        img = (image_map or {}).get(i)
        _add_content_slide(prs, sc, template_choice, insert_image, img)
    if insert_image:
        credits = {"title": "Credits", "content": "Images provided by Pexels: https://www.pexels.com"}
        _add_content_slide(prs, credits, template_choice, False, None)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def generate_presentation(
    user_text: str,
    number_of_slide: int,
    template_choice: str,
    presentation_title: str,
    presenter_name: str,
    insert_image: bool,
    language: str = "ko",
) -> tuple[bytes, list[dict[str, str]], str]:
    if OPENAI_API_KEY:
        raw = await chat_development(build_user_message(number_of_slide, user_text, language))
        slides_content = parse_response(raw)
        source = "openai"
    else:
        slides_content = builtin_slides(user_text, number_of_slide)
        source = "builtin"
    if not slides_content:
        slides_content = builtin_slides(user_text, number_of_slide)
    image_map: dict[int, bytes | None] = {}
    if insert_image and PEXELS_API_KEY:
        async with httpx.AsyncClient(timeout=30) as client:
            for i, sc in enumerate(slides_content):
                url = await search_pexels_images(sc.get("keyword", sc["title"]))
                if url:
                    r = await client.get(url)
                    image_map[i] = r.content if r.status_code == 200 else None
    pptx = create_ppt_bytes(
        slides_content,
        template_choice,
        presentation_title or user_text[:80],
        presenter_name,
        insert_image,
        image_map,
    )
    return pptx, slides_content, source
